from __future__ import annotations

import asyncio

from markitdown import MarkItDown

from src.services.document_database.ocr.doc_router import ProcessDecision
from src.services.document_database.ocr.page_format import format_document_with_pages
from src.services.document_database.ocr.types.interface import FileReader


class PptwReader(FileReader):
    async def read(
        self,
        filename: str,
        path: str,
        fate: ProcessDecision,
    ) -> str:
        if fate == ProcessDecision.OCR:
            # PPTX files don't require OCR
            raise ValueError("OCR is not applicable for PPTX files")

        file_path = self._resolve_file_path(path, filename)

        try:
            page_texts = await asyncio.to_thread(
                self._extract_slide_texts_sync, str(file_path)
            )
            return format_document_with_pages(filename, page_texts)
        except Exception:
            md = MarkItDown(enable_plugins=False)  # Set to True to enable plugins
            result = await asyncio.to_thread(md.convert, str(file_path))
            return format_document_with_pages(filename, [result.text_content])

    def _extract_slide_texts_sync(self, file_path: str) -> list[str]:
        from pptx import Presentation

        presentation = Presentation(file_path)
        page_texts: list[str] = []
        for slide in presentation.slides:
            lines: list[str] = []
            for shape in slide.shapes:
                shape_text = self._extract_shape_text(shape)
                if shape_text:
                    lines.append(shape_text)
            page_texts.append("\n".join(lines).strip())
        return page_texts or [""]

    def _extract_shape_text(self, shape) -> str:
        if getattr(shape, "has_text_frame", False):
            text = (shape.text or "").strip()
            if text:
                return text

        if getattr(shape, "has_table", False):
            table_lines: list[str] = []
            for row in shape.table.rows:
                row_values = [
                    cell.text.strip()
                    for cell in row.cells
                    if isinstance(cell.text, str) and cell.text.strip()
                ]
                if row_values:
                    table_lines.append(" | ".join(row_values))
            return "\n".join(table_lines).strip()

        return ""

    def get_page_count(self, filename: str, path: str) -> int:
        """
        Count slides in a PPTX file (each slide = 1 page).

        Uses python-pptx to count the number of slides in the presentation.
        Falls back to 1 if python-pptx is not installed or if the file cannot be read.
        """
        file_path = self._resolve_file_path(path, filename)
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))
            return len(prs.slides)
        except Exception:
            return 1
